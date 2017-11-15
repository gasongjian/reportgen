# -*- coding: utf-8 -*-
"""
Created on Tue Nov  8 20:05:36 2016
@author: JSong
"""
'''
pptx 用的单位是pptx.util.Emu,  英语单位
pptx.util.Inches(1)=914400
pptx.util.Pt(1)=12700
pptx.util.Cm(1)=360000
slide的大小:
pptx.Presentation().slide_height
pptx.Presentation().slide_width
'''

import os
import time


import pandas as pd
import numpy as np
#import matplotlib.pyplot as plt
import config


from pptx import Presentation
from pptx.chart.data import ChartData,XyChartData,BubbleChartData
from pptx.enum.chart import XL_CHART_TYPE
from pptx.util import Inches, Pt, Emu
from pptx.enum.chart import XL_LEGEND_POSITION
#from pptx.enum.chart import XL_LABEL_POSITION
from pptx.dml.color import RGBColor





chart_list={\
"AREA":[1,"ChartData"],\
"AREA_STACKED":[76,"ChartData"],\
"AREA_STACKED_100":[77,"ChartData"],\
"THREE_D_AREA":[-4098,"ChartData"],\
"THREE_D_AREA_STACKED":[78,"ChartData"],\
"THREE_D_AREA_STACKED_100":[79,"ChartData"],\
"BAR_CLUSTERED":[57,"ChartData"],\
"BAR_TWO_WAY":[57,"ChartData"],\
"BAR_OF_PIE":[71,"ChartData"],\
"BAR_STACKED":[58,"ChartData"],\
"BAR_STACKED_100":[59,"ChartData"],\
"THREE_D_BAR_CLUSTERED":[60,"ChartData"],\
"THREE_D_BAR_STACKED":[61,"ChartData"],\
"THREE_D_BAR_STACKED_100":[62,"ChartData"],\
"BUBBLE":[15,"BubbleChartData"],\
"BUBBLE_THREE_D_EFFECT":[87,"BubbleChartData"],\
"COLUMN_CLUSTERED":[51,"ChartData"],\
"COLUMN_STACKED":[52,"ChartData"],\
"COLUMN_STACKED_100":[53,"ChartData"],\
"THREE_D_COLUMN":[-4100,"ChartData"],\
"THREE_D_COLUMN_CLUSTERED":[54,"ChartData"],\
"THREE_D_COLUMN_STACKED":[55,"ChartData"],\
"THREE_D_COLUMN_STACKED_100":[56,"ChartData"],\
"CYLINDER_BAR_CLUSTERED":[95,"ChartData"],\
"CYLINDER_BAR_STACKED":[96,"ChartData"],\
"CYLINDER_BAR_STACKED_100":[97,"ChartData"],\
"CYLINDER_COL":[98,"ChartData"],\
"CYLINDER_COL_CLUSTERED":[92,"ChartData"],\
"CYLINDER_COL_STACKED":[93,"ChartData"],\
"CYLINDER_COL_STACKED_100":[94,"ChartData"],\
"DOUGHNUT":[-4120,"ChartData"],\
"DOUGHNUT_EXPLODED":[80,"ChartData"],\
"LINE":[4,"ChartData"],\
"LINE_MARKERS":[65,"ChartData"],\
"LINE_MARKERS_STACKED":[66,"ChartData"],\
"LINE_MARKERS_STACKED_100":[67,"ChartData"],\
"LINE_STACKED":[63,"ChartData"],\
"LINE_STACKED_100":[64,"ChartData"],\
"THREE_D_LINE":[-4101,"ChartData"],\
"PIE":[5,"ChartData"],\
"PIE_EXPLODED":[69,"ChartData"],\
"PIE_OF_PIE":[68,"ChartData"],\
"THREE_D_PIE":[-4102,"ChartData"],\
"THREE_D_PIE_EXPLODED":[70,"ChartData"],\
"PYRAMID_BAR_CLUSTERED":[109,"ChartData"],\
"PYRAMID_BAR_STACKED":[110,"ChartData"],\
"PYRAMID_BAR_STACKED_100":[111,"ChartData"],\
"PYRAMID_COL":[112,"ChartData"],\
"PYRAMID_COL_CLUSTERED":[106,"ChartData"],\
"PYRAMID_COL_STACKED":[107,"ChartData"],\
"PYRAMID_COL_STACKED_100":[108,"ChartData"],\
"RADAR":[-4151,"ChartData"],\
"RADAR_FILLED":[82,"ChartData"],\
"RADAR_MARKERS":[81,"ChartData"],\
"STOCK_HLC":[88,"ChartData"],\
"STOCK_OHLC":[89,"ChartData"],\
"STOCK_VHLC":[90,"ChartData"],\
"STOCK_VOHLC":[91,"ChartData"],\
"SURFACE":[83,"ChartData"],\
"SURFACE_TOP_VIEW":[85,"ChartData"],\
"SURFACE_TOP_VIEW_WIREFRAME":[86,"ChartData"],\
"SURFACE_WIREFRAME":[84,"ChartData"],\
"XY_SCATTER":[-4169,"XyChartData"],\
"XY_SCATTER_LINES":[74,"XyChartData"],\
"XY_SCATTER_LINES_NO_MARKERS":[75,"XyChartData"],\
"XY_SCATTER_SMOOTH":[72,"XyChartData"],\
"XY_SCATTER_SMOOTH_NO_MARKERS":[73,"XyChartData"]}


def df_to_table(slide,df,left,top,width,height,index_names=False,columns_names=True):
    '''将pandas数据框添加到slide上，并生成pptx上的表格
    输入：
    slide：PPT的一个页面，由pptx.Presentation().slides.add_slide()给定
    df：需要转换的数据框
    lef,top: 表格在slide中的位置
    width,height: 表格在slide中的大小
    index_names: Bool,是否需要显示行类别的名称
    columns_names: Bool,是否需要显示列类别的名称
    返回：
    返回带表格的slide
    '''
    df=pd.DataFrame(df)
    rows, cols = df.shape
    res = slide.shapes.add_table(rows+columns_names, cols+index_names, left, top, width, height)
    # 固定表格的宽度
    '''
    for c in range(cols+rownames):
        res.table.columns[c].width = colwidth
        res.table.rows[c].width = colwidth
    '''
    # Insert the column names
    if columns_names:
        for col_index, col_name in enumerate(list(df.columns)):
            cell=res.table.cell(0,col_index+index_names)
            #cell.text_frame.fit_text(max_size=12)
            #cell.text_frame.text='%s'%(col_name)
            cell.text = '%s'%(col_name)
    if index_names:
        for col_index, col_name in enumerate(list(df.index)):
            cell=res.table.cell(col_index+columns_names,0)
            cell.text = '%s'%(col_name)
            #cell.text_frame.fit_text(max_size=12)
    m = df.as_matrix()
    for row in range(rows):
        for col in range(cols):
            cell=res.table.cell(row+columns_names, col+index_names)
            cell.text = '%s'%(m[row, col])
            #cell.text_frame.fit_text(max_size=12)


def df_to_chartdata(df,datatype,number_format=None):
    '''
    根据给定的图表数据类型生成相应的数据
    Chartdata:一般的数据
    XyChartData: 散点图数据
    BubbleChartData:气泡图数据
    '''
    if isinstance(df,pd.Series):
        df=pd.DataFrame(df)
    df.fillna(0,inplace=True)
    datatype=datatype.lower()
    if datatype == 'chartdata':
        chart_data = ChartData()
        chart_data.categories = ['%s'%(c) for c in list(df.index)]
        for col_name in df.columns:
            chart_data.add_series('%s'%(col_name),list(df[col_name]),number_format)
        return chart_data
    if datatype == 'xychartdata':
        chart_data=XyChartData()
        if not isinstance(df,list):
            df=[df]
        for d in df:
            series_name='%s'%(d.columns[0])+' vs '+'%s'%(d.columns[1])
            series_ = chart_data.add_series(series_name)
            for i in range(len(d)):
                series_.add_data_point(d.iloc[i,0], d.iloc[i,1])
        return chart_data
    if datatype == 'bubblechartdata':
        chart_data=BubbleChartData()
        if not isinstance(df,list):
            df=[df]
        for d in df:
            series_name='%s'%(d.columns[0])+' vs '+'%s'%(d.columns[1])
            series_ = chart_data.add_series(series_name)
            for i in range(len(d)):
                series_.add_data_point(d.iloc[i,0],d.iloc[i,1],d.iloc[i,2])
        return chart_data



def plot_table(prs,df,layouts=[0,5],title=u'我是标题',summary=u'我是简短的结论',footnote=''):
    '''根据给定的数据，在给定的prs上新增一页表格ppt
    输入：
    prs: PPT文件接口
    df: 数据框
    layouts: [0]为PPT母版顺序，[1]为母版内的版式顺序
    输出：
    更新后的prs
    '''
    df=pd.DataFrame(df)
    slide_width=prs.slide_width
    slide_height=prs.slide_height
    # 可能需要修改以适应更多的情形
    title_only_slide = prs.slide_masters[layouts[0]].slide_layouts[layouts[1]]
    slide = prs.slides.add_slide(title_only_slide)
    #title=u'这里是标题'
    slide.shapes.title.text = title
    left,top = Emu(0.05*slide_width), Emu(0.10*slide_height)
    width,height = Emu(0.7*slide_width), Emu(0.1*slide_height)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    #summary=u'这里是一些简短的结论'
    txBox.text_frame.text=summary
    # 绘制表格
    '''添加自适应的表格大小
    默认最大12*6，width=0.80,height=0.70
    left=0.1,top=0.25
    '''
    R,C=df.shape
    width=max(0.5,min(1,C/6.0))*0.80
    height=max(0.5,min(1,R/12.0))*0.70
    left=0.5-width/2
    top=0.25
    left=Emu(left*slide_width)
    top=Emu(top*slide_height)
    width=Emu(width*slide_width)
    height=Emu(height*slide_height)
    df_to_table(slide,df,left,top,width,height,index_names=True)
    
        # 添加脚注 footnote=u'这里是脚注'
    if footnote:
        left,top = Emu(0.025*slide_width), Emu(0.95*slide_height)
        width,height = Emu(0.70*slide_width), Emu(0.10*slide_height)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        #p = text_frame.paragraphs[0]
        p=txBox.text_frame.paragraphs[0]
        p.text=footnote
        p.font.size = Pt(10)
        p.font.language_id = 3076
        p.font.name='Microsoft YaHei UI'
        p.font.color.rgb=RGBColor(127,127,127)
        try:
            txBox.text_frame.fit_text(max_size=10)
        except:
            pass
            #print('cannot fit the size of font')
    return prs


def plot_textbox(prs,texts,title=u'我是文本框页标题',summary=u'我是内容',footnote='',layouts=[0,0]):
    '''
    只绘制一个文本框，用于目录、小结等
    '''
    slide_width=prs.slide_width
    slide_height=prs.slide_height
    # 可能需要修改以适应更多的情形
    title_only_slide = prs.slide_masters[layouts[0]].slide_layouts[layouts[1]]
    slide = prs.slides.add_slide(title_only_slide)
    #title=u'这里是标题'
    slide.shapes.title.text = title
    # 绘制副标题
    if summary:
        left,top = Emu(0.15*slide_width), Emu(0.10*slide_height)
        width,height = Emu(0.7*slide_width), Emu(0.1*slide_height)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.text=summary
    # 绘制主体
    left,top = Emu(0.15*slide_width), Emu(0.20*slide_height)
    width,height = Emu(0.7*slide_width), Emu(0.7*slide_height)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text=texts

    # 添加脚注 footnote=u'这里是脚注'
    if footnote:
        left,top = Emu(0.025*slide_width), Emu(0.95*slide_height)
        width,height = Emu(0.70*slide_width), Emu(0.10*slide_height)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        #p = text_frame.paragraphs[0]
        p=txBox.text_frame.paragraphs[0]
        p.text=footnote
        p.font.size = Pt(10)
        p.font.language_id = 3076
        p.font.name='Microsoft YaHei UI'
        p.font.color.rgb=RGBColor(127,127,127)
        try:
            txBox.text_frame.fit_text(max_size=10)
        except:
            pass
            #print('cannot fit the size of font')
    return prs

def plot_picture(prs,img_path,layouts=[0,0],title=u'我是文本框页标题',summary='',\
footnote=''):
    '''
    只插入一张图片，用于目录、小结等
    '''
    slide_width=prs.slide_width
    slide_height=prs.slide_height
    # 可能需要修改以适应更多的情形
    title_only_slide = prs.slide_masters[layouts[0]].slide_layouts[layouts[1]]
    slide = prs.slides.add_slide(title_only_slide)
    #title=u'这里是标题'
    slide.shapes.title.text = title
    if summary:
        left,top = Emu(0.05*slide_width), Emu(0.10*slide_height)
        width,height = Emu(0.7*slide_width), Emu(0.1*slide_height)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        txBox.text_frame.text=summary
    left,top = Emu(0.15*slide_width), Emu(0.2*slide_height)
    height=Emu(0.7*slide_height)
    slide.shapes.add_picture(img_path, left, top, height=height)
    # 添加脚注 footnote=u'这里是脚注'
    if footnote:
        left,top = Emu(0.025*slide_width), Emu(0.95*slide_height)
        width,height = Emu(0.70*slide_width), Emu(0.10*slide_height)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        #p = text_frame.paragraphs[0]
        p=txBox.text_frame.paragraphs[0]
        p.text=footnote
        p.font.size = Pt(10)
        p.font.language_id = 3076
        p.font.name='Microsoft YaHei UI'
        p.font.color.rgb=RGBColor(127,127,127)
        try:
            txBox.text_frame.fit_text(max_size=10)
        except:
            pass
            #print('cannot fit the size of font')
    return prs



def plot_chart(prs,df,chart_type,title=u'我是标题',summary=u'我是简短的结论',\
footnote=None,chart_format=None,layouts=[0,0],has_data_labels=True):
    '''
    直接将数据绘制到一张ppt上，且高度定制化
    默认都有图例，且图例在下方
    默认都有数据标签
    '''

    slide_width=prs.slide_width
    slide_height=prs.slide_height
    # 可能需要修改以适应更多的情形
    # layouts[0]代表第几个母版，layouts[1]代表母版中的第几个版式
    title_only_slide = prs.slide_masters[layouts[0]].slide_layouts[layouts[1]]
    slide = prs.slides.add_slide(title_only_slide)
    # 添加标题 title=u'这里是标题'
    try:
        slide.shapes.title.text = title
    except:
        print('请检查模板,脚本没有找到合适的slide')
        return
    # 添加结论 summary=u'这里是一些简短的结论'
    #summary_loc=[0.10,0.14,0.80,0.15]
    left,top = Emu(config.summary_loc[0]*slide_width), Emu(config.summary_loc[1]*slide_height)
    width,height = Emu(config.summary_loc[2]*slide_width), Emu(config.summary_loc[3]*slide_height)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text=summary
    txBox.text_frame.paragraphs[0].font.language_id = 3076
    try:
        txBox.text_frame.fit_text(max_size=12)
    except:
        pass
        #print('cannot fit the size of font')


    # 添加脚注 footnote=u'这里是脚注'
    if footnote:
        left,top = Emu(0.025*slide_width), Emu(0.95*slide_height)
        width,height = Emu(0.70*slide_width), Emu(0.10*slide_height)
        txBox = slide.shapes.add_textbox(left, top, width, height)
        #p = text_frame.paragraphs[0]
        p=txBox.text_frame.paragraphs[0]
        p.text=footnote
        p.font.size = Pt(10)
        p.font.language_id = 3076
        p.font.name='Microsoft YaHei UI'
        p.font.color.rgb=RGBColor(127,127,127)
        try:
            txBox.text_frame.fit_text(max_size=10)
        except:
            pass
            #print('cannot fit the size of font')


    # 插入图表
    chart_type_code=chart_list[chart_type][1]
    chart_data=df_to_chartdata(df,chart_type_code)
    #left, top = Emu(0.05*slide_width), Emu(0.20*slide_height)
    #width, height = Emu(0.85*slide_width), Emu(0.70*slide_height)
    #chart_loc=[0.10,0.30,0.80,0.60]
    left, top = Emu(config.chart_loc[0]*slide_width), Emu(config.chart_loc[1]*slide_height)
    width, height = Emu(config.chart_loc[2]*slide_width), Emu(config.chart_loc[3]*slide_height)

    chart=slide.shapes.add_chart(chart_list[chart_type.upper()][0], \
    left, top, width, height, chart_data).chart

    if chart_type_code in [-4169,72,73,74,75]:
        return

    font_default_size=Pt(10)
    # 添加图例
    if (df.shape[1]>1) or (chart_type=='PIE'):
        chart.has_legend = True
        chart.legend.font.size=font_default_size
        chart.legend.position = XL_LEGEND_POSITION.BOTTOM
        chart.legend.include_in_layout = False

    try:
        chart.category_axis.tick_labels.font.size=font_default_size
    except:
        pass#暂时不知道怎么处理
    try:
        chart.value_axis.tick_labels.font.size=font_default_size
    except:
        pass
    # 添加数据标签

    non_available_list=['BUBBLE','BUBBLE_THREE_D_EFFECT','XY_SCATTER',\
    'XY_SCATTER_LINES','PIE']
    # 大致检测是否采用百分比
    # 1、单选题每列的和肯定是100，顶多相差+-5
    # 2、多选题每一列的和大于100，但单个的小于100.此处可能会有误判，但暂时无解
    # 3、可能会有某一列全为0，此时单独考虑
    if  ((df.sum()[df.sum()!=0]>90).all()) and ((df<=100).all().all()) and (u'总体' not in df.index):
        # 数据条的数据标签格式
        #number_format1='0.0"%"'
        number_format1=config.number_format_chart_1
        # 坐标轴的数据标签格式
        #number_format2='0"%"'
        number_format2=config.number_format_chart_label
    else:
        number_format1='0.00'
        number_format2='0.0'

    if (chart_type not in non_available_list) or (chart_type == 'PIE'):
        plot = chart.plots[0]
        plot.has_data_labels = True
        plot.data_labels.font.size = font_default_size
        plot.data_labels.number_format = number_format1
        #plot.data_labels.number_format_is_linked=True
        #data_labels = plot.data_labels
        #plot.data_labels.position = XL_LABEL_POSITION.BEST_FIT
    if (chart_type not in non_available_list):
        #chart.value_axis.maximum_scale = 1
        if df.shape[1]==1:
            chart.value_axis.has_major_gridlines = False
        else:
            chart.value_axis.has_major_gridlines = True
        tick_labels = chart.value_axis.tick_labels
        tick_labels.number_format = number_format2
        tick_labels.font.size = font_default_size

    # 修改纵坐标格式
    '''
    tick_labels = chart.value_axis.tick_labels
    tick_labels.number_format = '0"%"'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(10)
    '''

    # 填充系列的颜色
    ''' 最好的方法还是修改母版文件中的主题颜色，这里只提供方法
    if df.shape[1]==1:
        chart.series[0].fill()
    '''

    # 自定义format
    if chart_format:
        for k in chart_format:
            exec('chart.'+k+'='+'%s'%(chart_format[k]))

    return prs

    '''
    if chart_type == 'BAR_TWO_WAY':
        chart
    '''


def plot_cover(prs,title=u'reportgen工具包封面',layouts=[0,0],xspace=8,yspace=6):
    try:
        from delaunay import Delaunay2D
    except:
        return 
    slide_width=prs.slide_width
    slide_height=prs.slide_height
    # 可能需要修改以适应更多的情形
    title_only_slide = prs.slide_masters[layouts[0]].slide_layouts[layouts[1]]
    slide = prs.slides.add_slide(title_only_slide)

    ## 随机生成连接点
    seeds=np.round(np.dot(np.random.rand((xspace-1)*(yspace-1),2),np.diag([slide_width,slide_height])))
    # 添加左边点
    tmp=np.linspace(0,slide_height,yspace)
    seeds=np.concatenate((seeds,np.array([[0]*len(tmp),tmp]).T))
    # 添加上边点
    tmp=np.linspace(0,slide_width,xspace)[1:]
    seeds=np.concatenate((seeds,np.array([tmp,[0]*len(tmp)]).T))
    # 添加右边点
    tmp=np.linspace(0,slide_height,yspace)[1:]
    seeds=np.concatenate((seeds,np.array([[slide_width]*len(tmp),tmp]).T))
    # 添加下边点
    tmp=np.linspace(0,slide_width,xspace)[1:-1]
    seeds=np.concatenate((seeds,np.array([tmp,[slide_height]*len(tmp)]).T))

    # 构造三角剖分，生成相应的三角形和平面图数据
    center = np.mean(seeds, axis=0)
    t=np.sqrt(slide_width**2+slide_height**2)/2
    dt = Delaunay2D(center, 2**(np.floor(np.log2(t))+1))
    for s in seeds:
        dt.AddPoint(s)
    tri=dt.exportTriangles()
    graph=np.zeros((len(seeds),len(seeds)))
    for t in tri:
        graph[t[0],t[1]]=1
        graph[t[1],t[2]]=1
        graph[t[0],t[2]]=1
        graph[t[1],t[0]]=1
        graph[t[2],t[1]]=1
        graph[t[2],t[1]]=1


    from pptx.enum.shapes import MSO_CONNECTOR
    from pptx.enum.shapes import MSO_SHAPE
    shapes = slide.shapes
    # 添加连接线
    for i in range(len(seeds)):
        for j in range(len(seeds)):
            if (i<j) and graph[i,j]==1:
                shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT, Emu(seeds[i,0]), Emu(seeds[i,1]), Emu(seeds[j,0]), Emu(seeds[j,1]))
    # 添加圆点，原点的半径符合高斯分布
    radius=slide_width/100
    for i in range(len(seeds)):
        eps=np.random.normal(scale=radius*0.2)
        left=Emu(seeds[i,0])-radius-eps
        top=Emu(seeds[i,1])-radius-eps
        width=height=2*(radius+eps)
        shape=shapes.add_shape(
        MSO_SHAPE.OVAL,left, top, width, height)
        shape.line.width=Emu(0)
        fill = shape.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(218,227,243)

    # 添加标题
    left,top = Emu(0), Emu(0.4*slide_height)
    width,height = Emu(1*slide_width), Emu(0.2*slide_height)
    shape=shapes.add_shape(
    MSO_SHAPE.RECTANGLE,left, top, width, height)
    shape.line.width=Emu(0)
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = RGBColor(0,176,240)
    shape.text=title

    # 添加脚注
    left,top = Emu(0.72*slide_width), Emu(0.93*slide_height)
    width,height = Emu(0.25*slide_width), Emu(0.07*slide_height)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text='POWERED BY REPORTGEN'

    # 添加LOGO
    import os
    logo_path=os.path.join(os.path.split(__file__)[0],'logo.png')
    if os.path.exists(logo_path):
        left,top = Emu(0.65*slide_width), Emu(0.94*slide_height)
        height=Emu(0.06*slide_height)
        slide.shapes.add_picture(logo_path, left, top, height=height)
    return prs


def slides_data_gen(slides_data,chart_type_default='COLUMN_CLUSTERED'):
    '''自动补全pptx数据信息
    slides_data: 默认需可迭代
    每一页PPT设定为四个元素：标题、结论、主题、脚注
    return
    ------
    slides_data: 每一页ppt所需要的元素[
        {title:,#标题
        summary:,#结论
        data:,# DataFrame数据、文本数据、图片地址等
        slide_type:,#chart、table、text
        chart_type:图表类型
        data_config:,#字典格式，绘制data其他所需要的相关参数，保留字段，暂时不用
        footnote:,#脚注
        layouts:#该slide使用的ppt版式
        },]
    filename: 缺省以时间命名
    template:使用的模板
    '''


    title=''
    summary=''
    footnote=''
    # 处理slides_data数据
    if (not isinstance(slides_data,list)) and (not isinstance(slides_data,tuple)):
        slides_data=[slides_data]
    # 自动计算图表的格式
    #np.issubdtype(a.as_matrix().dtype,np.number)
    # 补全相关信息
    slides_data_new=[]
    for i in range(len(slides_data)):
        slide=slides_data[i]
        # 补全相关信息,数据处理部分待定
        if not isinstance(slide,dict):
            slide={'data':slide}
            slide['title']=title
            slide['summary']=summary
            slide['footnote']=footnote
            slide['layouts']='auto'
            slide['data_config']=None
            if isinstance(slide['data'],pd.core.frame.DataFrame):
                slide['slide_type']='chart'
                slide['chart_type']=chart_type_default
            elif isinstance(slide['data'],pd.core.series.Series):
                slide['data']=pd.DataFrame(slide['data'])
                slide['slide_type']='chart'
                slide['chart_type']=chart_type_default
            elif isinstance(slide['data'],str) and os.path.exists(slide['data']):
                slide['slide_type']='picture'
                slide['chart_type']=None
            elif isinstance(slide['data'],str) and not(os.path.exists(slide['data'])):
                slide['slide_type']='textbox'
                slide['chart_type']=''
            else:
                print('未知的数据格式，请检查数据')
                slide['slide_type']=None
                slide['chart_type']=None
                continue
        elif isinstance(slide,dict):
            if 'data' not in slide:
                print('没有找到需要的数据，请检查')
                slide['slide_type']=None
                slide['chart_type']=None
                continue
            if isinstance(slide['data'],pd.core.series.Series):
                slide['data']=pd.DataFrame(slide['data'])
            if 'title' not in slide:
                slide['title']=title
            if 'summary' not in slide:
                slide['summary']=summary
            if 'footnote' not in slide:
                slide['footnote']=footnote
            if 'layouts' not in slide:
                slide['layouts']='auto'
            if 'data_config' not in slide:
                slide['data_config']=None
            slide['chart_type']=None if 'chart_type' not in slide else slide['chart_type']
            if 'slide_type' not in slide:
                if isinstance(slide['data'],pd.core.frame.DataFrame):
                    slide['slide_type']='chart'
                    slide['chart_type']=chart_type_default
                elif isinstance(slide['data'],str) and os.path.exists(slide['data']):
                    print('test')
                    slide['slide_type']='picture'
                    slide['chart_type']=''
                elif isinstance(slide['data'],str) and not(os.path.exists(slide['data'])):
                    slide['slide_type']='textbox'
                    slide['chart_type']=''
                else:
                    print('未知的数据格式，请检查数据')
                    slide['slide_type']=None
                    slide['chart_type']=None
                    continue
        slides_data_new.append(slide)

    return slides_data_new





class Report():
    '''
    报告自动生成工具
    r=Report(filename='')
    r.add_cover(title='reportgen')
    r.add_slides([])
    r.save()
    '''
    def __init__(self,filename=None,chart_type_default='COLUMN_CLUSTERED'):

        #self.template=template
        self.filename=filename
        self.chart_type_default=chart_type_default
        if filename is None:
            if os.path.exists('template.pptx'):
                prs=Presentation('template.pptx')
            else:
                prs=Presentation()
        else :
            prs=Presentation(filename)
        self.prs=prs
        title_only_slide=self._layouts()
        if title_only_slide:
            layouts=title_only_slide[0]
        else:
            layouts=[0,0]
        self.layouts_default=layouts


    def _layouts(self):
        '''给定pptx文件，自动识别标题等版式
        '''
        slide_width=self.prs.slide_width
        slide_height=self.prs.slide_height
        title_only_slide=[]
        #blank_slide=[]
        for i in range(len(self.prs.slide_masters)):
            slides=self.prs.slide_masters[i]
            #print('第{}个有{}个版式'.format(i,len(slides.slide_layouts)))
            for j in range(len(slides.slide_layouts)):
                slide=slides.slide_layouts[j]
                title_slide=0
                placeholder_size=0
                for k in range(len(slide.shapes)):
                    shape=slide.shapes[k]
                    if shape.is_placeholder and shape.has_text_frame:
                        left,top=shape.left/slide_width,shape.top/slide_height
                        height=shape.height/slide_height
                        if left<1 and top<1 and height<1 and left>0 and top>0 and height>0:
                            placeholder_size+=1
                        #print('left={:.2f},top={:.2f},height={:.2f}'.format(left,top,height))
                        if left<0.15 and top<0.15 and height <0.25:
                            title_slide+=1
                #print('{}个文本占位符,{}个title'.format(placeholder_size,title_slide))
                if placeholder_size==1 and title_slide==1:
                    title_only_slide.append([i,j])
                #if placeholder_size==0:
                    #blank_slide.append((i,j))s
        return title_only_slide



    def get_texts(self):
        # one for each text run in presentation
        text_runs = []
    
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if not shape.has_text_frame:
                    continue
                for paragraph in shape.text_frame.paragraphs:
                    for run in paragraph.runs:
                        text_runs.append(run.text)
        return text_runs
        
    def get_images(self):
        try:
            from PIL import Image as PIL_Image
            from io import BytesIO
        except:
            print('please install the PIL.')
            return
        if not os.path.exists('.\\images'):
            os.mkdir('.\\images')
        for slide in self.prs.slides:
            for shape in slide.shapes:
                if 'Image' in str(type(shape)) or 'Picture' in str(type(shape)):
                    shape_image=shape.image
                    #filename='.\\images\\'+shape_image.filename
                    r=str(np.random.randint(99)).zfill(2)
                    filename='.\\images\\images'+time.strftime('_%m%d%H%M%S', time.localtime())+r+'.'+shape_image.ext
                    p = PIL_Image.open(BytesIO(shape_image.blob))
                    p.save(filename)
                    #print('save {}'.format(shape_image.filename))




    def add_slides(self,slides_data,chart_type_default=None):
        '''
        slides_data: 每一页ppt所需要的元素[
            {title:,#标题
            summary:,#结论
            data:,# DataFrame数据、文本数据、图片地址等
            slide_type:,#chart、table、text
            chart_type:图表类型
            data_config:,#字典格式，绘制data其他所需要的相关参数，保留字段，暂时不用
            footnote:,#脚注
            layouts:#该slide使用的ppt版式
            },]
        '''
        if chart_type_default is None:
            chart_type_default=self.chart_type_default
        slides_data=slides_data_gen(slides_data,chart_type_default)
        for slide in slides_data:
            slide_type=slide['slide_type']
            title=slide['title']
            summary=slide['summary']
            footnote=slide['footnote']
            layouts=self.layouts_default if slide['layouts'] == 'auto' else slide['layouts']
            data=slide['data']
            chart_type=slide['chart_type'] if 'chart_type' in slide else None
            #data_config=slide['data_config']#暂时没有用该参数
            if (slide_type is None) or (not isinstance(slide_type,str)):
                continue
            if slide_type == 'chart':
                self.prs=plot_chart(self.prs,data,chart_type=chart_type,layouts=layouts,\
                title=title,summary=summary,footnote=footnote);
            elif slide_type == 'table':
                self.prs=plot_table(self.prs,data,layouts=layouts,title=title,summary=summary,\
                footnote=footnote);
            elif slide_type in ['textbox','text']:
                self.prs=plot_textbox(self.prs,data,layouts=layouts,title=title,summary=summary,\
                footnote=footnote);
            elif slide_type in ['picture','figure']:
                self.prs=plot_picture(self.prs,data,layouts=layouts,title=title,summary=summary,\
                footnote=footnote);
        return self


    def add_cover(self,title='',author='',style='default',layouts='auto',size=[8,6]):
        title =u'reportgen自动生成工具' if len(title)==0 else title
        #author =u'report' if len(author)==0 else author
        layouts=self.layouts_default if layouts == 'auto' else layouts
        if style == 'default':
            self.prs=plot_cover(self.prs,title=title,layouts=layouts,xspace=size[0],yspace=size[1]);
        return self

        
    def location_suggest(self,num=1,rate=0.78):
        '''统一管理slides各个模块的位置
        parameter
        --------
        num: 主体内容（如图、外链图片、文本框等）的个数，默认从左到右依次排列
        rate: 主体内容的宽度综合
        
        return
        -----
        locations: dict格式. l代表left,t代表top,w代表width，h代表height
        '''
        slide_width,slide_height=self.prs.slide_width,self.prs.slide_height
        if 'summary_loc' in config.__dict__:
            summary_loc=config.summary_loc
        else:
            summary_loc=[0.10,0.14,0.80,0.15]
            
        if 'footnote_loc' in config.__dict__:
            footnote_loc=config.footnote_loc
        else:
            footnote_loc=[0.025,0.95,0.70,0.06]           

        if 'data_loc' in config.__dict__:
            data_loc=config.data_loc
        else:
            data_loc=[0.11,0.30,0.78,0.60]
            
        locations={}
        locations['summary']={'l':Emu(summary_loc[0]*slide_width),'t':Emu(summary_loc[1]*slide_height),\
                 'w':Emu(summary_loc[2]*slide_width),'h':Emu(summary_loc[3]*slide_height)}
        
        locations['footnote']={'l':Emu(footnote_loc[0]*slide_width),'t':Emu(footnote_loc[1]*slide_height),\
                 'w':Emu(footnote_loc[2]*slide_width),'h':Emu(footnote_loc[3]*slide_height)}
        # 主体部分只有一个的情形
        '''
        控制主体的宽度为78%，且居中显示。
        '''
        if num>1:
            left=[(1-rate)*(i+1)/(float(num)+1)+rate*i/float(num) for i in range(num)]
            top=[data_loc[1]]*num
            width=[rate/float(num)]*num
            height=[data_loc[3]]*num 
            locations['data']=[{'l':Emu(left[i]*slide_width),'t':Emu(top[i]*slide_height),\
                     'w':Emu(width[i]*slide_width),'h':Emu(height[i]*slide_height)} for i in range(num)]
        else:
            locations['data']=[{'l':Emu(data_loc[0]*slide_width),'t':Emu(data_loc[1]*slide_height),\
                     'w':Emu(data_loc[2]*slide_width),'h':Emu(data_loc[3]*slide_height)}]

        return locations


    def add_slide(self,data=[],title='',summary='',footnote='',layouts='auto',**kwarg):
        '''通用格式
        data=[{'data':,'slide_type':,'type':,},] # 三个是必须字段，其他根据slide_type不同而不同
        number_format_data: 图的数据标签
        number_format_tick: 横纵坐标的数据标签
        
        '''
        #slide_width=self.prs.slide_width
        #slide_height=self.prs.slide_height
        
        # 标准化data格式
        if not(isinstance(data,list)):
            data=[data]
        for i,d in enumerate(data):
            if not(isinstance(d,dict)):               
                if isinstance(d,(pd.core.frame.DataFrame,pd.core.frame.Series)):
                    slide_type='chart'
                    chart_type=self.chart_type_default
                elif isinstance(d,str) and os.path.exists(d):
                    slide_type='picture'
                    chart_type=''
                elif isinstance(d,str) and not(os.path.exists(d)):
                    slide_type='textbox'
                    chart_type=''
                else:
                    print('未知的数据格式，请检查数据')
                    slide_type=''
                    chart_type=''
                data[i]={'data':d,'slide_type':slide_type,'type':chart_type}

        # 各个模板的位置
        locations=self.location_suggest(len(data))
        summary_loc=locations['summary']
        footnote_loc=locations['footnote']
        data_loc=locations['data']
        
        # 选取的板式
        if layouts == 'auto':
            layouts=self.layouts_default
        title_only_slide = self.prs.slide_masters[layouts[0]].slide_layouts[layouts[1]]
        slide = self.prs.slides.add_slide(title_only_slide)
        
        #输出标题
        slide.shapes.title.text = title
      
        # 输出副标题 summary
        if summary:
            txBox = slide.shapes.add_textbox(summary_loc['l'], summary_loc['t'], summary_loc['w'], summary_loc['h'])
            txBox.text_frame.text=summary
            txBox.text_frame.paragraphs[0].font.language_id = 3076
            try:
                txBox.text_frame.fit_text(max_size=12)
            except:
                pass

    
        # 输出脚注 footnote
        if footnote:
            txBox = slide.shapes.add_textbox(footnote_loc['l'], footnote_loc['t'], footnote_loc['w'], footnote_loc['h'])
            #p = text_frame.paragraphs[0]
            p=txBox.text_frame.paragraphs[0]
            p.text=footnote
            p.font.size = Pt(10)
            p.font.language_id = 3076
            p.font.name='Microsoft YaHei UI'
            p.font.color.rgb=RGBColor(127,127,127)
            try:
                txBox.text_frame.fit_text(max_size=10)
            except:
                pass
                #print('cannot fit the size of font')
                
        for i,dd in  enumerate(data):
            slide_type=dd['slide_type']
            left,top=data_loc[i]['l'],data_loc[i]['t']
            width,height=data_loc[i]['w'],data_loc[i]['h']
            chart_type=dd['type'] if 'type' in dd else self.chart_type_default
            if slide_type in ['table']:
                # 绘制表格
                '''针对表格大小修正
                R,C=dd['data'].shape
                width=max(0.5,min(1,C/6.0))*width
                height=max(0.5,min(1,R/12.0))*height
                left=0.5-width/2
                top=0.25
                '''
                df_to_table(slide,dd['data'],left,top,width,height,index_names=True)
            elif slide_type in ['textbox']:
                # 输出文本框
                txBox = slide.shapes.add_textbox(left, top, width, height)
                txBox.text_frame.text=dd['data']
                txBox.text_frame.paragraphs[0].font.language_id = 3076
                try:
                    txBox.text_frame.fit_text(max_size=12)
                except:
                    pass
            elif slide_type in ['picture','figure']:
                slide.shapes.add_picture(dd['data'], left, top, height=height)
            elif slide_type in ['chart']:
                # 插入图表
                chart_type_code=chart_list[chart_type][1]
                chart_data=df_to_chartdata(dd['data'],chart_type_code)          
                chart=slide.shapes.add_chart(chart_list[chart_type.upper()][0],left, top, width, height, chart_data).chart
            
                if chart_type_code in [-4169,72,73,74,75]:
                    continue                   
                font_default_size=Pt(10) if 'font_default_size' not in config.__dict__ else config.font_default_size
                # 添加图例
                if (dd['data'].shape[1]>1) or (chart_type=='PIE'):
                    chart.has_legend = True
                    chart.legend.font.size=font_default_size
                    chart.legend.position = XL_LEGEND_POSITION.BOTTOM
                    chart.legend.include_in_layout = False          
                try:
                    chart.category_axis.tick_labels.font.size=font_default_size
                except:
                    pass#暂时不知道怎么处理
                try:
                    chart.value_axis.tick_labels.font.size=font_default_size
                except:
                    pass
                # 添加数据标签
            
                non_available_list=['BUBBLE','BUBBLE_THREE_D_EFFECT','XY_SCATTER','XY_SCATTER_LINES','PIE']
                # 数据标签数值格式
                number_format1='0.00' if 'number_format_data' not in dd else dd['number_format_data']
                number_format2='0.0' if 'number_format_tick' not in dd else dd['number_format_tick']          
                if (chart_type not in non_available_list) or (chart_type == 'PIE'):
                    plot = chart.plots[0]
                    plot.has_data_labels = True
                    plot.data_labels.font.size = font_default_size
                    plot.data_labels.number_format = number_format1
                    #data_labels = plot.data_labels
                    #plot.data_labels.position = XL_LABEL_POSITION.BEST_FIT
                if (chart_type not in non_available_list):
                    #chart.value_axis.maximum_scale = 1
                    if dd['data'].shape[1]==1:
                        chart.value_axis.has_major_gridlines = False
                    else:
                        chart.value_axis.has_major_gridlines = True
                    tick_labels = chart.value_axis.tick_labels
                    tick_labels.number_format = number_format2
                    tick_labels.font.size = font_default_size                
                 
        return self


    def save(self,filename=None):
        filename=self.filename+time.strftime('_%Y%m%d%H%M.pptx', time.localtime()) if filename is None else filename
        self.prs.save(filename)





if __name__ == '__main__':
    '''
    暂时用于保存一些信息
    '''
    prs = Presentation()
    slide_width=prs.slide_width
    slide_height=prs.slide_height

    '''
    模板文件放在,可自己修改，其中不能增加板式
    C:\Anaconda2\Lib\site-packages\pptx\templates\default.pptx
    复杂用法：
    prs = Presentation()
    '''

    #标题页面
    title_slide = prs.slide_layouts[0]
    #重点列表
    bullet_slide = prs.slide_layouts[1]
    # 只有一个标题
    title_only_slide = prs.slide_layouts[5]
    # 空页面
    blank_slide = prs.slide_layouts[6]

    # 第一页：标题
    slide = prs.slides.add_slide(title_slide)
    title = slide.shapes.title
    subtitle = slide.placeholders[1]
    title.text = u"半自动化报告生成"
    subtitle.text = "powered by python-pptx"

    # 第二页：文本框
    slide = prs.slides.add_slide(blank_slide)
    left = top = width = height = Inches(1)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    tf = txBox.text_frame
    tf.text = "This is text inside a textbox"
    p = tf.add_paragraph()
    p.text = "This is a second paragraph that's bold"
    p.font.bold = True
    p = tf.add_paragraph()
    p.text = "This is a third paragraph that's big"
    p.font.size = Pt(40)

    # 第三页：表格
    '''
    默认表格样式
    '''
    slide = prs.slides.add_slide(title_only_slide)
    shapes = slide.shapes
    shapes.title.text = u'表格示例'
    #d=pd.read_csv('data.csv',encoding='gbk')
    #t=pd.crosstab(d['Q2'],d['Q3'])
    left=Emu(0.15*slide_width)
    top=Emu(0.3*slide_height)
    width=Emu(0.7*slide_width)
    height=Emu(0.6*slide_height)
    #df_to_table(slide,t,left,top,width,height,rownames=True)

    # 第4、5、6、7页：图表

    t1=pd.DataFrame({u'满意度':[2,4,3,3.5],u'重要度':[2.2,3.5,2.6,3.8]},columns=[u'满意度',u'重要度'])
    t2=pd.DataFrame({u'满意度':[2,4,3,3.5],u'重要度':[3.2,4.5,3.6,4.8]},columns=[u'满意度',u'重要度'])
    t3=pd.DataFrame({u'满意度':[2,4,3,3.5],u'重要度':[2.2,3.5,2.6,3.8],\
    u'大小':[3,2,6,8]},columns=[u'满意度',u'重要度',u'大小'])
    t1.rename(index={0:u'张',1:u'王',2:u'李',3:u'宋'},inplace=True)
    plot_chart(prs,t1,'COLUMN_CLUSTERED')
    plot_chart(prs,t1,'COLUMN_STACKED')
    plot_chart(prs,t1,'BAR_CLUSTERED')
    #plot_chart(prs,t,'BAR_CLUSTERED')
    plot_chart(prs,[t1,t2],'XY_SCATTER')
    plot_chart(prs,t3,'BUBBLE')
    plot_chart(prs,pd.DataFrame({u'满意度':[2,4,3,3.5]}),'PIE')

    # 第8页：自定义图
    slide = prs.slides.add_slide(title_only_slide)
    slide.shapes.title.text = u'高度定制化'
    left,top = Emu(0.12*slide_width), Emu(0.2*slide_height)
    width,height = Emu(0.7*slide_width), Emu(0.1*slide_height)
    txBox = slide.shapes.add_textbox(left, top, width, height)
    txBox.text_frame.text=u'这里是一些简短的结论'
    # define chart data ---------------------
    '''
    chart_data = ChartData()
    chart_data.categories = ['East', 'West', 'Midwest']
    chart_data.add_series('Series 1', (19.2, 21.4, 16.7))
    chart_data.add_series('Series 2', (23, 16, 19))
    '''
    chart_data=df_to_chartdata(t1,'ChartData')
    # add chart to slide --------------------
    x, y = Emu(0.05*slide_width), Emu(0.35*slide_height)
    cx, cy = Emu(0.76*slide_width), Emu(0.55*slide_height)
    chart=slide.shapes.add_chart(XL_CHART_TYPE.BAR_CLUSTERED, \
    x, y, cx, cy, chart_data).chart



    '''
    # ----样式-------
    chart.chart_style: 样式(ppt自带的1到48)
    # ----坐标轴-------
    axis=chart.category_axis: X轴坐标控制
    axis=chart.value_axis   :Y周坐标控制
    axis.visible: 坐标轴是否可见
    axis.has_major_gridlines：添加主要网格线
    axis.major_gridlines: 设置主要网格线
    axis.has_minor_gridlines: 添加次要网格线
    axis.minor_gridlines: 设置次要网格线
    axis.major_tick_mark: 主要刻度线类型(XL_TICK_MARK.OUTSIDE: 无，内部，外部，交叉)
    axis.minor_tick_mark: 主要刻度线类型(XL_TICK_MARK.OUTSIDE: 无，内部，外部，交叉)
    axis.maximum_scale:  最大值
    axis.minimum_scale：最小值
    axis.tick_label_position: 坐标轴标签
    axis.tick_labels.font:  字体设置(共8个维度,如.bold=True, .size=Pt(12))
    axis.tick_labels.number_format: 数字格式('0"%"')
    axis.tick_labels.number_format_is_linked
    axis.tick_labels.offset
    # ----数据标签-------
    plot = chart.plots[0]
    plot.has_data_labels = True
    data_labels = plot.data_labels
    data_labels.font.size = Pt(13)
    # from pptx.dml.color import RGBColor
    data_labels.font.color.rgb = RGBColor(0x0A, 0x42, 0x80)
    # from pptx.enum.chart import XL_LABEL_POSITION
    data_labels.position = XL_LABEL_POSITION.INSIDE_END
    # ----图例-------
    # from pptx.enum.chart import XL_LEGEND_POSITION
    chart.has_legend = True
    chart.legend.position = XL_LEGEND_POSITION.RIGHT
    chart.legend.include_in_layout = False
    '''
    '''
    category_axis = chart.category_axis
    category_axis.has_major_gridlines = True
    category_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    category_axis.tick_labels.font.italic = True
    category_axis.tick_labels.font.size = Pt(24)
    value_axis = chart.value_axis
    value_axis.maximum_scale = 50.0
    value_axis.minor_tick_mark = XL_TICK_MARK.OUTSIDE
    value_axis.has_minor_gridlines = True
    tick_labels = value_axis.tick_labels
    tick_labels.number_format = '0"%"'
    tick_labels.font.bold = True
    tick_labels.font.size = Pt(14)
    '''

    # 添加致谢页
    #slide = prs.slides.add_slide(prs.slide_layouts[8])
    prs.save('test.pptx')
